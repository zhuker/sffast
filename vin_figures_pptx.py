#!/usr/bin/env python3
"""
vin_figures_pptx.py - Export applicable figures for a VIN to a PowerPoint file.

Creates one slide per figure page with:
  - Header: Category, Figure ID+Page, Description
  - PNG image of the figure
  - Callouts with part numbers and details
  - Cross-references and bulletin references

Bulletin images are added as separate slides at the end.

Usage: .venv/bin/python vin_figures_pptx.py [VIN]
       Default VIN: JF1GD70655L510047 (2005 STI)
"""

import io
import os
import sys
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from sffastus_parser import is_valid_subaru_vin
from sffastus_database import SffastDatabase, PartMatch


# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Image: 1280x640 -> 2:1 aspect ratio
IMG_WIDTH = Inches(12.0)
IMG_HEIGHT = Inches(6.0)  # 12/2 = 6
IMG_LEFT = (SLIDE_WIDTH - IMG_WIDTH) // 2
IMG_TOP = Inches(0.6)

# Text areas
HEADER_TOP = Inches(0.05)
HEADER_HEIGHT = Inches(0.5)
CALLOUT_TOP = IMG_TOP + IMG_HEIGHT + Inches(0.05)
CALLOUT_LEFT = Inches(0.3)
CALLOUT_WIDTH = SLIDE_WIDTH - Inches(0.6)

FONT_NAME = "Courier New"
FONT_SIZE = Pt(6)
HEADER_FONT_SIZE = Pt(12)


def format_callout_line(callout_code, variant, part):
    """Format a single part match line (like vin_figures.py display_part)."""
    cc = callout_code
    if part.variant:
        cc = cc + "*" + part.variant

    extra = []
    if part.usage_notes:
        extra.append(part.usage_notes)
    if part.part_spec:
        extra.append(part.part_spec)
    extra_str = f"[{', '.join(extra)}]" if extra else ""

    chain_str = ""
    if part.itca_chain:
        parts_str = []
        for link in part.itca_chain:
            if link.bulletin_figure:
                parts_str.append(
                    f"{link.part_number}({link.code.label}"
                    f" FIG {link.bulletin_figure}-{link.bulletin_page}"
                    f" {link.bulletin_label})")
            else:
                parts_str.append(f"{link.part_number}({link.code.label})")
        chain_str = " -> " + " -> ".join(parts_str)

    if part.bulletin_figure:
        bltn = f"FIG {part.bulletin_figure}-{part.bulletin_page} {part.bulletin_label}"
        if chain_str:
            chain_str = f" [{bltn}]{chain_str}"
        else:
            chain_str = f" [{bltn}]"

    return cc, extra_str, chain_str


def build_callout_text(db, model_rec, fig, page, vehicle):
    """Build callout text lines for a figure page (mirrors vin_figures.py logic)."""
    lines = []

    callouts = db.get_fig_callouts(model_rec, fig, page, vehicle=vehicle)
    pg_callouts = sorted([c for c in callouts if c.source == 'part_group'],
                         key=lambda c: c.callout_code)
    inv_callouts = sorted([c for c in callouts if c.source == 'inventory'],
                          key=lambda c: c.callout_code)
    pg_bases = set(c.callout_code for c in pg_callouts)

    for c in pg_callouts:
        if c.parts:
            p0 = c.parts[0]
            cc0, extra0, chain0 = format_callout_line(c.callout_code, c.variant, p0)
            buf = f"{cc0:8s}   {p0.part_number:14s} {p0.description} {extra0} {chain0}"

            alt_parts = []
            for pN in c.parts[1:]:
                ccN, extraN, chainN = format_callout_line(c.callout_code, c.variant, pN)
                pieces = []
                if ccN != cc0:
                    pieces.append(ccN)
                if pN.part_number != p0.part_number:
                    pieces.append(pN.part_number)
                if pN.description != p0.description:
                    pieces.append(pN.description)
                if extraN != extra0:
                    pieces.append(extraN)
                if chainN != chain0:
                    pieces.append(chainN)
                alt_parts.append(" ".join(pieces))

            if alt_parts:
                buf += " OR " + " OR ".join(alt_parts)
            lines.append(buf)
        else:
            cc = c.callout_code
            if c.variant:
                cc = cc + "*" + c.variant
            lines.append(f"{cc:8s}   {'--':14s} {c.description}")

    for c in inv_callouts:
        if c.callout_code in pg_bases:
            continue
        if not c.parts:
            continue
        pn = c.parts[0].part_number
        lines.append(f"{c.callout_code:8s}   {pn:14s} {c.description}")

    return lines, callouts


def build_xref_text(db, model_rec, fig, page):
    """Build cross-reference text lines."""
    lines = []
    xrefs = db.get_fig_xrefs(model_rec, fig, page)
    for xr in xrefs:
        xr_info = db.get_fig_info(model_rec, xr.ref_figure)
        xr_desc = xr_info.desc_en if xr_info else db.get_figname(xr.ref_figure)
        xr_cat = ""
        if xr_info:
            cat184 = db.get_fig_category(model_rec, xr_info.fig_group_code)
            if cat184:
                xr_cat = f" [{cat184.desc_en}]"
        lines.append(f"  -> FIG {xr.ref_figure}{xr_cat} {xr_desc}")
    return lines


def build_bulletin_ref_lines(callouts):
    """Extract bulletin references from callout parts."""
    bulletin_refs = set()
    for c in callouts:
        for p in c.parts:
            if p.bulletin_figure:
                bulletin_refs.add((p.bulletin_figure, p.bulletin_page, p.bulletin_label))
            if p.itca_chain:
                for link in p.itca_chain:
                    if link.bulletin_figure:
                        bulletin_refs.add((link.bulletin_figure, link.bulletin_page, link.bulletin_label))
    lines = []
    for bf, bp, bl in sorted(bulletin_refs):
        lines.append(f"  -> FIG {bf}-{bp} {bl}")
    return lines


def get_fig_png(db, model_rec, fig, page):
    """Extract figure image as PNG bytes."""
    wand_img = db.get_fig_img(model_rec, fig, page)
    if wand_img is None:
        return None
    with wand_img:
        return wand_img.make_blob('png')


def add_text_frame(slide, left, top, width, height, text_lines, font_size=FONT_SIZE,
                   font_name=FONT_NAME, bold=False, color=None):
    """Add a text frame with monospace text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        if bold:
            p.font.bold = True
        if color:
            p.font.color.rgb = color
        p.space_before = Pt(0)
        p.space_after = Pt(0)

    return txBox


def add_header(slide, text):
    """Add a header text box at the top of the slide."""
    add_text_frame(slide, Inches(0.3), HEADER_TOP, SLIDE_WIDTH - Inches(0.6),
                   HEADER_HEIGHT, [text], font_size=HEADER_FONT_SIZE,
                   font_name="Arial", bold=True)


def add_image(slide, png_bytes):
    """Add a PNG image to the slide."""
    stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(stream, IMG_LEFT, IMG_TOP, IMG_WIDTH, IMG_HEIGHT)


def get_header_text(db, model_rec, fig, page, prefix=""):
    """Build header string: [Category] FIG fig-page Description (label)"""
    fig183 = db.get_fig_info(model_rec, fig)
    fig_desc = fig183.desc_en if fig183 else db.get_figname(fig)

    cat_str = ""
    if fig183:
        cat184 = db.get_fig_category(model_rec, fig183.fig_group_code)
        if cat184:
            cat_str = f"[{cat184.desc_en}] "

    fig89 = db.get_fig_page(model_rec, fig, page)
    label = ""
    if fig89 and fig89.label:
        label = " ".join(fig89.label.split())
    label_str = f" ({label})" if label else ""

    return f"{cat_str}{prefix}FIG {fig}-{page} {fig_desc}{label_str}"


def main():
    vin = sys.argv[1] if len(sys.argv) > 1 else "JF1GD70655L510047"

    if not is_valid_subaru_vin(vin):
        print(f"Error: '{vin}' is not a valid Subaru VIN")
        sys.exit(1)

    print(f"VIN: {vin}")

    with SffastDatabase.open() as db:
        print("Resolving VIN...")
        try:
            vehicle = db.resolve_vin(vin)
        except LookupError as e:
            print(str(e))
            sys.exit(1)

        model_rec = vehicle.model_rec
        spec = vehicle.spec

        print(f"  Model: {vehicle.vin_rec.model_code}")
        if spec:
            print(f"  Applied Model: {spec.applied_model} / {spec.engine} / {spec.transmission}")

        # Get applicable pages
        all_pages = db.get_vin_figures(vehicle)
        applicable_figs = [p for p in all_pages if p.type == 'figure']
        bulletin_pages = [p for p in all_pages if p.type == 'bulletin']

        print(f"  Applicable figures: {len(applicable_figs)}")
        print(f"  I&S bulletins: {len(bulletin_pages)}")

        # Get bulletin matches for back-references
        bulletin_matches = db.get_vin_bulletins(vehicle)
        # Group by (bulletin_figure, bulletin_page)
        bulletin_refs_by_page = defaultdict(list)
        for bm in bulletin_matches:
            bulletin_refs_by_page[(bm.figure, bm.page)].append(bm)

        # Group figures by category
        by_fig_page = defaultdict(list)
        for p in applicable_figs:
            by_fig_page[(p.figure, p.page)].append(p)

        category_figures = defaultdict(list)
        uncategorized = []
        for fig, page in sorted(by_fig_page.keys()):
            fig183 = db.get_fig_info(model_rec, fig)
            if fig183:
                category_figures[fig183.fig_group_code].append((fig, page))
            else:
                uncategorized.append((fig, page))

        sorted_cats = sorted(category_figures.keys())

        # Create presentation
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]  # blank slide

        # Title slide
        slide = prs.slides.add_slide(blank_layout)
        title_lines = [
            f"Subaru Parts Catalog",
            f"VIN: {vin}",
            f"Model: {vehicle.vin_rec.model_code}",
        ]
        if spec:
            title_lines.extend([
                f"Applied Model: {spec.applied_model}",
                f"Engine: {spec.engine}  Trans: {spec.transmission}  Trim: {spec.trim_level}",
                f"Body: {spec.body_config}  Drivetrain: {spec.drivetrain}",
            ])
        title_lines.append(f"Production: {vehicle.vehicle_date}")
        title_lines.append(f"")
        title_lines.append(f"Figures: {len(by_fig_page)}  Bulletins: {len(bulletin_pages)}")
        add_text_frame(slide, Inches(1), Inches(1.5), Inches(11), Inches(5),
                       title_lines, font_size=Pt(18), font_name="Arial", bold=True)

        # Process figure slides
        total_figures = len(by_fig_page)
        fig_count = 0

        # Collect all bulletin refs seen across all figure slides
        all_bulletin_figs_seen = set()

        def add_figure_slide(fig, page):
            nonlocal fig_count
            fig_count += 1
            print(f"  [{fig_count}/{total_figures}] FIG {fig}-{page}")

            slide = prs.slides.add_slide(blank_layout)

            # Header
            header = get_header_text(db, model_rec, fig, page)
            add_header(slide, header)

            # Image
            png_bytes = get_fig_png(db, model_rec, fig, page)
            if png_bytes:
                add_image(slide, png_bytes)

            # Callouts
            callout_lines, callouts = build_callout_text(db, model_rec, fig, page, vehicle)

            # Cross-references
            xref_lines = build_xref_text(db, model_rec, fig, page)

            # Bulletin references from parts
            bref_lines = build_bulletin_ref_lines(callouts)

            all_lines = callout_lines + xref_lines + bref_lines

            if all_lines:
                # Calculate available height below image
                available_height = SLIDE_HEIGHT - CALLOUT_TOP - Inches(0.1)
                add_text_frame(slide, CALLOUT_LEFT, CALLOUT_TOP,
                               CALLOUT_WIDTH, available_height, all_lines)

            # Track bulletin refs for later
            for c in callouts:
                for p in c.parts:
                    if p.bulletin_figure:
                        all_bulletin_figs_seen.add((p.bulletin_figure, p.bulletin_page))
                    if p.itca_chain:
                        for link in p.itca_chain:
                            if link.bulletin_figure:
                                all_bulletin_figs_seen.add((link.bulletin_figure, link.bulletin_page))

        print("\nGenerating figure slides...")
        for cat_code in sorted_cats:
            for fig, page in category_figures[cat_code]:
                add_figure_slide(fig, page)

        for fig, page in uncategorized:
            add_figure_slide(fig, page)

        # Bulletin slides at the end
        # Collect unique bulletin pages from both get_vin_figures and parts references
        bulletin_set = set()
        for b in bulletin_pages:
            bulletin_set.add((b.figure, b.page))
        bulletin_set.update(all_bulletin_figs_seen)
        sorted_bulletins = sorted(bulletin_set)

        if sorted_bulletins:
            print(f"\nGenerating {len(sorted_bulletins)} bulletin slides...")
            for i, (bfig, bpage) in enumerate(sorted_bulletins, 1):
                print(f"  [{i}/{len(sorted_bulletins)}] Bulletin FIG {bfig}-{bpage}")

                slide = prs.slides.add_slide(blank_layout)

                # Header
                header = get_header_text(db, model_rec, bfig, bpage, prefix="I&S Bulletin ")
                add_header(slide, header)

                # Image
                png_bytes = get_fig_png(db, model_rec, bfig, bpage)
                if png_bytes:
                    add_image(slide, png_bytes)

                # Build referencing text: which figures/parts reference this bulletin
                ref_lines = []
                refs = bulletin_refs_by_page.get((bfig, bpage), [])
                for bm in refs:
                    itca_label = bm.itca_code.label if bm.itca_code else "direct"
                    ref_lines.append(
                        f"FIG {bm.source_figure}-{bm.source_page}"
                        f"  callout {bm.source_callout}"
                        f"  part {bm.source_part}"
                        f"  ({itca_label})")

                if ref_lines:
                    header_line = f"Referenced by {len(ref_lines)} part(s):"
                    all_lines = [header_line] + ref_lines
                    available_height = SLIDE_HEIGHT - CALLOUT_TOP - Inches(0.1)
                    add_text_frame(slide, CALLOUT_LEFT, CALLOUT_TOP,
                                   CALLOUT_WIDTH, available_height, all_lines)

        # Save
        os.makedirs("output", exist_ok=True)
        out_path = f"output/{vin}.pptx"
        prs.save(out_path)
        print(f"\nSaved: {out_path}")
        print(f"Total slides: {len(prs.slides)} (1 title + {fig_count} figures + {len(sorted_bulletins)} bulletins)")


if __name__ == "__main__":
    main()
